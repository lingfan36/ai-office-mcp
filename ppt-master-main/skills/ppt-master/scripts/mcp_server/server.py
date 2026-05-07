#!/usr/bin/env python3
"""PPT Master MCP Server — 26 tools: inspect, edit, animate, design, render, export."""
import json, sys, copy
from pathlib import Path

_HERE       = Path(__file__).resolve().parent
_SCRIPTS    = _HERE.parent
_RENDERERS  = _SCRIPTS / "renderers"
_PROJ_ROOT  = _SCRIPTS.parent.parent.parent
_OUT_DIR    = _PROJ_ROOT / "projects"
_TMPL_DIR   = _HERE / "pptx_templates"
_TMPL_DIR.mkdir(exist_ok=True)

sys.path.insert(0, str(_RENDERERS))
sys.path.insert(0, str(_SCRIPTS))

from mcp.server.fastmcp import FastMCP
mcp = FastMCP("ppt-master")

# ── helpers ────────────────────────────────────────────────────────────────────
PX = 9525  # 1 px @ 96 DPI in EMU
_P  = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A  = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R  = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_P14= "http://schemas.microsoft.com/office/powerpoint/2010/main"

def _load(path: str):
    from pptx import Presentation
    p = Path(path)
    if not p.exists(): raise FileNotFoundError(f"Not found: {path}")
    return Presentation(str(p))

def _save(prs, path: str) -> str:
    p = Path(path); p.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(p)); return str(p)

def _emu(px):
    from pptx.util import Emu; return Emu(int(px * PX))

def _pxv(emu): return int(emu / PX)

def _rgb(h: str):
    from pptx.dml.color import RGBColor
    h = h.lstrip("#"); return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:6],16))

def _sdict(s, i: int) -> dict:
    d = {"index": i, "shape_id": s.shape_id, "name": s.name,
         "type": str(s.shape_type).split(".")[-1],
         "x": _pxv(s.left), "y": _pxv(s.top), "w": _pxv(s.width), "h": _pxv(s.height)}
    if s.has_text_frame: d["text"] = s.text_frame.text
    return d

def _ok(**kw): return json.dumps({"status": "ok", **kw}, ensure_ascii=False)
def _err(m):   return json.dumps({"status": "error", "message": m}, ensure_ascii=False)

def _parse_pptx_xml(xml_str: str):
    """Parse a PPTX XML fragment that uses p:/a:/r: prefixes without declarations."""
    from lxml import etree
    wrapped = (f'<root xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}" xmlns:p14="{_P14}">'
               f'{xml_str.strip()}</root>')
    return etree.fromstring(wrapped.encode("utf-8"))[0]

def _rId_of_slide(xml_slides, idx: int) -> str:
    from pptx.oxml.ns import qn
    return xml_slides[idx].get(qn("r:id")) or ""


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 1 — INSPECT / READ
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def get_presentation_info(pptx_path: str) -> str:
    """Return high-level info about a PPTX: slide count, canvas size, notes presence.

    Args:
        pptx_path: Absolute path to the PPTX file.
    Returns:
        JSON: status, slide_count, width_px, height_px, slides_with_notes (list of indices).
    """
    try:
        prs = _load(pptx_path)
        return _ok(
            slide_count=len(prs.slides),
            width_px=_pxv(prs.slide_width),
            height_px=_pxv(prs.slide_height),
            slides_with_notes=[i for i, s in enumerate(prs.slides) if s.has_notes_slide],
        )
    except Exception as e: return _err(str(e))


@mcp.tool()
def list_slides(pptx_path: str) -> str:
    """List all slides with index, title text, shape count, and notes flag.

    Args:
        pptx_path: Absolute path to the PPTX file.
    Returns:
        JSON: status, slides — array of {index, title, shape_count, has_notes}.
    """
    try:
        prs = _load(pptx_path)
        out = []
        for i, slide in enumerate(prs.slides):
            title = next(
                (s.text_frame.text for s in slide.shapes
                 if s.has_text_frame and s.name.lower().startswith("title")), ""
            ) or next(
                (s.text_frame.text[:80] for s in slide.shapes if s.has_text_frame), ""
            )
            out.append({"index": i, "title": title,
                        "shape_count": len(slide.shapes), "has_notes": slide.has_notes_slide})
        return _ok(slides=out)
    except Exception as e: return _err(str(e))


@mcp.tool()
def get_slide_content(pptx_path: str, slide_index: int) -> str:
    """Return all shapes on a slide: position, size, type, and text.

    Args:
        pptx_path: Absolute path to the PPTX file.
        slide_index: Zero-based slide index.
    Returns:
        JSON: status, slide_index, shapes — array of shape dicts.
        Each shape: index, shape_id, name, type, x, y, w, h, text (if text frame).
    """
    try:
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range (0–{len(prs.slides)-1})")
        slide = prs.slides[slide_index]
        return _ok(slide_index=slide_index,
                   shapes=[_sdict(s, i) for i, s in enumerate(slide.shapes)])
    except Exception as e: return _err(str(e))


@mcp.tool()
def extract_all_text(pptx_path: str) -> str:
    """Extract all text from a PPTX, grouped by slide.

    Args:
        pptx_path: Absolute path to the PPTX file.
    Returns:
        JSON: status, slides — array of {index, texts: [str]}.
    """
    try:
        prs = _load(pptx_path)
        return _ok(slides=[
            {"index": i, "texts": [s.text_frame.text.strip()
                                   for s in slide.shapes
                                   if s.has_text_frame and s.text_frame.text.strip()]}
            for i, slide in enumerate(prs.slides)
        ])
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 2 — SLIDE MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def add_blank_slide(pptx_path: str, position: int = -1) -> str:
    """Add a blank slide. -1 (default) appends at end, otherwise inserts at position.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        position: Zero-based insert index. -1 = append.
    Returns:
        JSON: status, slide_count, new_slide_index.
    """
    try:
        prs = _load(pptx_path)
        prs.slides.add_slide(prs.slide_layouts[6])
        xml_slides = prs.slides._sldIdLst
        if 0 <= position < len(xml_slides) - 1:
            el = xml_slides[-1]
            xml_slides.remove(el)
            xml_slides.insert(position, el)
            new_idx = position
        else:
            new_idx = len(prs.slides) - 1
        _save(prs, pptx_path)
        return _ok(slide_count=len(prs.slides), new_slide_index=new_idx)
    except Exception as e: return _err(str(e))


@mcp.tool()
def remove_slide(pptx_path: str, slide_index: int) -> str:
    """Remove a slide from the presentation.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based index of slide to remove.
    Returns:
        JSON: status, slide_count.
    """
    try:
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        xml_slides = prs.slides._sldIdLst
        rId = _rId_of_slide(xml_slides, slide_index)
        if rId: prs.part.drop_rel(rId)
        xml_slides.remove(xml_slides[slide_index])
        _save(prs, pptx_path)
        return _ok(slide_count=len(prs.slides))
    except Exception as e: return _err(str(e))


@mcp.tool()
def reorder_slides(pptx_path: str, from_index: int, to_index: int) -> str:
    """Move a slide to a new position.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        from_index: Current zero-based index of the slide.
        to_index: Target zero-based index (slide inserted before this position).
    Returns:
        JSON: status, slide_count.
    """
    try:
        prs = _load(pptx_path)
        if not (0 <= from_index < len(prs.slides)):
            return _err(f"from_index {from_index} out of range")
        xml_slides = prs.slides._sldIdLst
        el = xml_slides[from_index]
        xml_slides.remove(el)
        xml_slides.insert(max(0, min(to_index, len(xml_slides))), el)
        _save(prs, pptx_path)
        return _ok(slide_count=len(prs.slides))
    except Exception as e: return _err(str(e))


@mcp.tool()
def duplicate_slide(pptx_path: str, slide_index: int) -> str:
    """Duplicate a slide, appending the copy at the end.

    Text and drawing shapes are copied. Embedded images may not transfer
    if their relationship parts cannot be resolved; re-add them manually.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based index of slide to copy.
    Returns:
        JSON: status, slide_count, new_slide_index.
    """
    try:
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        tmpl = prs.slides[slide_index]
        new_slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Replace spTree
        new_tree = new_slide.shapes._spTree
        for ch in list(new_tree): new_tree.remove(ch)
        for ch in tmpl.shapes._spTree: new_tree.append(copy.deepcopy(ch))
        # Copy background element if present
        bg = tmpl._element.find(f"{{{_P}}}cSld/{{{_P}}}bg")
        if bg is not None:
            cSld = new_slide._element.find(f"{{{_P}}}cSld")
            existing_bg = cSld.find(f"{{{_P}}}bg")
            if existing_bg is not None: cSld.remove(existing_bg)
            spTree = cSld.find(f"{{{_P}}}spTree")
            cSld.insert(list(cSld).index(spTree), copy.deepcopy(bg))
        new_idx = len(prs.slides) - 1
        _save(prs, pptx_path)
        return _ok(slide_count=len(prs.slides), new_slide_index=new_idx)
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 3 — CONTENT EDITING
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def update_text(pptx_path: str, slide_index: int, shape_index: int, new_text: str,
                font_size: float = 0, bold: bool = False, color: str = "",
                align: str = "") -> str:
    """Update the text of a shape. Uses \\n for paragraph breaks.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        shape_index: Zero-based shape index on that slide.
        new_text: Replacement text (\\n = paragraph break).
        font_size: Points. 0 = keep existing.
        bold: Apply bold when font_size > 0 or color is set.
        color: Hex color e.g. 'FF0000'. '' = keep existing.
        align: 'left' | 'center' | 'right' | '' (keep).
    Returns: JSON: status.
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    _AL = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    try:
        prs = _load(pptx_path)
        shape = prs.slides[slide_index].shapes[shape_index]
        if not shape.has_text_frame: return _err("Shape has no text frame")
        tf = shape.text_frame
        tf.text = new_text
        for para in tf.paragraphs:
            if align in _AL: para.alignment = _AL[align]
            for run in para.runs:
                if font_size > 0: run.font.size = Pt(font_size); run.font.bold = bold
                if color: run.font.color.rgb = _rgb(color)
        _save(prs, pptx_path); return _ok()
    except Exception as e: return _err(str(e))


@mcp.tool()
def delete_shape(pptx_path: str, slide_index: int, shape_index: int) -> str:
    """Remove a shape from a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        shape_index: Zero-based shape index.
    Returns: JSON: status, remaining_shapes.
    """
    try:
        prs = _load(pptx_path)
        slide = prs.slides[slide_index]
        sp = slide.shapes[shape_index]._element
        sp.getparent().remove(sp)
        _save(prs, pptx_path)
        return _ok(remaining_shapes=len(slide.shapes))
    except Exception as e: return _err(str(e))


@mcp.tool()
def update_shape_style(pptx_path: str, slide_index: int, shape_index: int,
                       fill_color: str = "", border_color: str = "",
                       border_width_pt: float = 0, font_color: str = "",
                       font_size: float = 0, bold: bool = False) -> str:
    """Change fill, border, and font style of a shape.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index / shape_index: Zero-based indices.
        fill_color: Hex fill e.g. 'E0F2FE'. '' = no change.
        border_color: Hex border. '' = no change.
        border_width_pt: Border width in pt. 0 = no change.
        font_color: Hex text color. '' = no change.
        font_size: Points. 0 = no change.
        bold: Applied when font_size > 0 or font_color is set.
    Returns: JSON: status.
    """
    from pptx.util import Pt
    try:
        prs = _load(pptx_path)
        shape = prs.slides[slide_index].shapes[shape_index]
        if fill_color: shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(fill_color)
        if border_color: shape.line.color.rgb = _rgb(border_color)
        if border_width_pt > 0: shape.line.width = Pt(border_width_pt)
        if (font_color or font_size > 0) and shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if font_size > 0: run.font.size = Pt(font_size); run.font.bold = bold
                    if font_color: run.font.color.rgb = _rgb(font_color)
        _save(prs, pptx_path); return _ok()
    except Exception as e: return _err(str(e))


@mcp.tool()
def add_text_box(pptx_path: str, slide_index: int,
                 x: float, y: float, w: float, h: float,
                 text: str, font_size: float = 24, color: str = "1E293B",
                 bold: bool = False, align: str = "left", font: str = "") -> str:
    """Add a formatted text box to a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        x, y: Top-left in pixels (96 DPI). w, h: Size in pixels.
        text: Content (\\n = paragraph break).
        font_size: Points (default 24). color: Hex (default '1E293B').
        bold: Bold text. align: 'left'|'center'|'right'.
        font: Font name. '' = theme default.
    Returns: JSON: status, shape_index.
    """
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    _AL = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    try:
        prs = _load(pptx_path)
        slide = prs.slides[slide_index]
        tb = slide.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.text = text
        for para in tf.paragraphs:
            para.alignment = _AL.get(align, PP_ALIGN.LEFT)
            for run in para.runs:
                run.font.size = Pt(font_size); run.font.bold = bold
                run.font.color.rgb = _rgb(color)
                if font: run.font.name = font
        _save(prs, pptx_path)
        return _ok(shape_index=len(slide.shapes) - 1)
    except Exception as e: return _err(str(e))


@mcp.tool()
def add_shape(pptx_path: str, slide_index: int, shape_type: str,
              x: float, y: float, w: float, h: float,
              fill_color: str = "0EA5E9", border_color: str = "",
              border_width_pt: float = 0) -> str:
    """Add an auto-shape to a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        shape_type: 'rectangle' | 'ellipse' | 'rounded_rectangle' | 'triangle' |
                    'diamond' | 'chevron' | 'pentagon' | 'hexagon' |
                    'right_arrow' | 'star_5'.
        x, y, w, h: Position and size in pixels.
        fill_color: Hex fill (default '0EA5E9').
        border_color: Hex border. '' = no border.
        border_width_pt: Border width in pt. 0 = hairline.
    Returns: JSON: status, shape_index, shape_id.
    """
    from pptx.util import Pt
    _SMAP = {"rectangle": 1, "rounded_rectangle": 5, "triangle": 7,
             "ellipse": 9, "hexagon": 10, "right_arrow": 13, "diamond": 4,
             "chevron": 52, "pentagon": 56, "star_5": 92}
    try:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        prs = _load(pptx_path)
        slide = prs.slides[slide_index]
        enum_val = MSO_AUTO_SHAPE_TYPE(_SMAP.get(shape_type.lower(), 1))
        shape = slide.shapes.add_shape(enum_val, _emu(x), _emu(y), _emu(w), _emu(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = _rgb(fill_color)
        if border_color:
            shape.line.color.rgb = _rgb(border_color)
            if border_width_pt > 0: shape.line.width = Pt(border_width_pt)
        else:
            shape.line.fill.background()
        _save(prs, pptx_path)
        return _ok(shape_index=len(slide.shapes) - 1, shape_id=shape.shape_id)
    except Exception as e: return _err(str(e))


@mcp.tool()
def add_image_to_slide(pptx_path: str, slide_index: int, image_source: str,
                       x: float, y: float, w: float, h: float) -> str:
    """Insert an image (local path or HTTPS URL) onto a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        image_source: Absolute local file path OR https:// URL.
        x, y, w, h: Position and size in pixels.
    Returns: JSON: status, shape_index.
    """
    import io
    try:
        prs = _load(pptx_path)
        slide = prs.slides[slide_index]
        if image_source.startswith("http"):
            import requests
            r = requests.get(image_source, headers={"User-Agent": "Mozilla/5.0"}, timeout=20)
            r.raise_for_status()
            src = io.BytesIO(r.content)
        else:
            src = image_source
        slide.shapes.add_picture(src, _emu(x), _emu(y), _emu(w), _emu(h))
        _save(prs, pptx_path)
        return _ok(shape_index=len(slide.shapes) - 1)
    except Exception as e: return _err(str(e))


@mcp.tool()
def add_table_to_slide(pptx_path: str, slide_index: int,
                       headers: list, rows: list,
                       x: float, y: float, w: float, h: float,
                       header_fill: str = "1E293B", header_font_color: str = "FFFFFF",
                       row_fill: str = "F1F5F9", alt_row_fill: str = "FFFFFF",
                       font_size: float = 14) -> str:
    """Insert a formatted table onto a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        headers: List of column header strings.
        rows: List of rows; each row is a list of cell value strings.
        x, y, w, h: Position and size in pixels.
        header_fill / header_font_color: Header style hex colors.
        row_fill / alt_row_fill: Alternating row background hex colors.
        font_size: Points (default 14).
    Returns: JSON: status, shape_index.
    """
    from pptx.util import Pt
    try:
        prs = _load(pptx_path)
        slide = prs.slides[slide_index]
        n_cols, n_rows = len(headers), len(rows) + 1
        tbl = slide.shapes.add_table(n_rows, n_cols, _emu(x), _emu(y), _emu(w), _emu(h)).table
        col_w = int(_emu(w) / n_cols)
        for col in tbl.columns: col.width = col_w

        def _fc(cell, text, bg, fg, bd=False):
            cell.fill.solid(); cell.fill.fore_color.rgb = _rgb(bg)
            cell.text_frame.text = str(text)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = _rgb(fg)
                    run.font.bold = bd

        for c, h in enumerate(headers):
            _fc(tbl.cell(0, c), h, header_fill, header_font_color, bd=True)
        for r, row_data in enumerate(rows):
            fill = row_fill if r % 2 == 0 else alt_row_fill
            for c, val in enumerate(row_data[:n_cols]):
                _fc(tbl.cell(r + 1, c), val, fill, "1E293B")

        _save(prs, pptx_path)
        return _ok(shape_index=len(slide.shapes) - 1)
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 4 — ANIMATIONS & TRANSITIONS
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def set_slide_transition(pptx_path: str, slide_index: int,
                         transition_type: str = "fade", duration_ms: int = 500,
                         advance_after_ms: int = 0, apply_to_all: bool = False) -> str:
    """Set the transition effect for one slide or all slides.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index (ignored when apply_to_all=True).
        transition_type: fade | push | wipe | split | strips | cover | random.
        duration_ms: Transition duration in ms (default 500).
        advance_after_ms: Auto-advance after N ms. 0 = manual (default).
        apply_to_all: Apply this transition to every slide.
    Returns: JSON: status, applied_to (list of indices).
    """
    try:
        from pptx_animations import create_transition_xml
        from lxml import etree
        prs = _load(pptx_path)
        targets = list(range(len(prs.slides))) if apply_to_all else [slide_index]
        if not apply_to_all and not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        trans_xml = create_transition_xml(
            transition_type, duration=duration_ms / 1000,
            advance_after=advance_after_ms / 1000 if advance_after_ms > 0 else None)
        trans_el = _parse_pptx_xml(trans_xml)
        tag = f"{{{_P}}}transition"
        for idx in targets:
            sp_el = prs.slides[idx]._element
            for old in sp_el.findall(tag): sp_el.remove(old)
            sp_el.append(copy.deepcopy(trans_el))
        _save(prs, pptx_path)
        return _ok(applied_to=targets)
    except Exception as e: return _err(str(e))


@mcp.tool()
def add_shape_animation(pptx_path: str, slide_index: int, shape_indices: list,
                        animation_type: str = "fade", trigger: str = "after-previous",
                        duration_ms: int = 500, delay_ms: int = 0) -> str:
    """Add entrance animations to shapes on a slide. Replaces any existing timing.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        shape_indices: Ordered list of zero-based shape indices to animate.
        animation_type: appear|fade|fly|cut|zoom|wipe|split|blinds|checkerboard|
            dissolve|random_bars|peek|wheel|box|circle|diamond|plus|strips|
            wedge|stretch|expand|swivel | 'mixed' | 'random'.
        trigger: 'after-previous' (auto-cascade) | 'on-click' | 'with-previous'.
        duration_ms: Per-shape duration in ms (default 500).
        delay_ms: Gap between sequential animations in ms (default 0).
    Returns: JSON: status, animated_shape_ids.
    """
    try:
        from pptx_animations import create_sequence_timing_xml, pick_animation_effect
        from lxml import etree
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        slide = prs.slides[slide_index]
        targets, ids = [], []
        for i, sidx in enumerate(shape_indices):
            if 0 <= sidx < len(slide.shapes):
                sid = slide.shapes[sidx].shape_id
                targets.append((sid, delay_ms, pick_animation_effect(animation_type, i)))
                ids.append(sid)
        if not targets: return _err("No valid shapes found")
        timing_xml = create_sequence_timing_xml(targets, duration=duration_ms / 1000,
                                                trigger=trigger)
        timing_el = _parse_pptx_xml(timing_xml)
        sp_el = slide._element
        for old in sp_el.findall(f"{{{_P}}}timing"): sp_el.remove(old)
        sp_el.append(timing_el)
        _save(prs, pptx_path)
        return _ok(animated_shape_ids=ids)
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 5 — DESIGN / THEME
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def set_slide_background(pptx_path: str, slide_index: int,
                         color: str = "", gradient_start: str = "",
                         gradient_end: str = "", gradient_angle: int = 135) -> str:
    """Set a slide's background to a solid color or two-stop linear gradient.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        color: Hex solid fill. Takes priority over gradient.
        gradient_start / gradient_end: Hex start/end colors for gradient.
        gradient_angle: Degrees (default 135 = top-left→bottom-right diagonal).
    Returns: JSON: status.
    """
    try:
        from lxml import etree
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        slide = prs.slides[slide_index]
        if color:
            fill = slide.background.fill
            fill.solid(); fill.fore_color.rgb = _rgb(color)
        elif gradient_start and gradient_end:
            ang = int((gradient_angle % 360) * 60000)
            gs = gradient_start.lstrip("#"); ge = gradient_end.lstrip("#")
            grad_xml = (f'<a:gradFill xmlns:a="{_A}" rotWithShape="1"><a:gsLst>'
                        f'<a:gs pos="0"><a:srgbClr val="{gs}"/></a:gs>'
                        f'<a:gs pos="100000"><a:srgbClr val="{ge}"/></a:gs>'
                        f'</a:gsLst><a:lin ang="{ang}" scaled="0"/></a:gradFill>')
            grad_el = etree.fromstring(grad_xml.encode())
            cSld = slide._element.find(f"{{{_P}}}cSld")
            bg = cSld.find(f"{{{_P}}}bg")
            if bg is None:
                bg = etree.SubElement(cSld, f"{{{_P}}}bg")
                spTree = cSld.find(f"{{{_P}}}spTree")
                if spTree is not None:
                    cSld.remove(bg); cSld.insert(list(cSld).index(spTree), bg)
            bgPr = bg.find(f"{{{_P}}}bgPr")
            if bgPr is None: bgPr = etree.SubElement(bg, f"{{{_P}}}bgPr")
            for ch in list(bgPr): bgPr.remove(ch)
            bgPr.append(grad_el)
            stretch = etree.SubElement(bgPr, f"{{{_A}}}stretch")
            etree.SubElement(stretch, f"{{{_A}}}fillRect")
        else:
            return _err("Provide 'color' or both 'gradient_start' and 'gradient_end'")
        _save(prs, pptx_path); return _ok()
    except Exception as e: return _err(str(e))


@mcp.tool()
def set_speaker_notes(pptx_path: str, slide_index: int, notes_text: str) -> str:
    """Set or replace the speaker notes on a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        notes_text: Notes content. \\n for line breaks.
    Returns: JSON: status.
    """
    try:
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        prs.slides[slide_index].notes_slide.notes_text_frame.text = notes_text
        _save(prs, pptx_path); return _ok()
    except Exception as e: return _err(str(e))


@mcp.tool()
def get_speaker_notes(pptx_path: str, slide_index: int) -> str:
    """Read the speaker notes from a slide.

    Args:
        pptx_path: Path to PPTX.
        slide_index: Zero-based slide index.
    Returns: JSON: status, notes (string, empty if none).
    """
    try:
        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        slide = prs.slides[slide_index]
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        return _ok(notes=notes)
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 6 — SMART RENDERING
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def render_slide(pptx_path: str, slide_spec_json: str, slide_index: int = -1) -> str:
    """Add or replace a single slide using the premium layout engine.

    Args:
        pptx_path: Path to an existing PPTX (or non-existent path to create new).
        slide_spec_json: JSON object for one slide, same format as render_presentation's
            slides array entry. Must include 'layout' field.
            Example: {"layout":"cover","title":"Hello","page":"01/05"}
        slide_index: Zero-based index to replace. -1 (default) appends.
    Returns: JSON: status, slide_count, affected_index, pptx_path.
    """
    try:
        spec = json.loads(slide_spec_json)
    except json.JSONDecodeError as e:
        return _err(f"Invalid JSON: {e}")
    try:
        import engine, theme, layouts
        from pptx import Presentation
        from pptx.util import Emu
        from pptx.oxml.ns import qn
        p = Path(pptx_path)
        if p.exists():
            prs = _load(pptx_path)
        else:
            prs = Presentation()
            prs.slide_width  = Emu(theme.SLIDE_W_PX * theme.PX_TO_EMU)
            prs.slide_height = Emu(theme.SLIDE_H_PX * theme.PX_TO_EMU)
        fn = engine.DISPATCH.get(spec.get("layout"))
        if fn is None:
            return _err(f"Unknown layout '{spec.get('layout')}'. Call list_layouts() for options.")
        new_slide = prs.slides.add_slide(prs.slide_layouts[6])
        fn(new_slide, spec)
        xml_slides = prs.slides._sldIdLst
        n = len(xml_slides)
        if 0 <= slide_index < n - 1:
            old_el = xml_slides[slide_index]
            old_rId = _rId_of_slide(xml_slides, slide_index)
            if old_rId: prs.part.drop_rel(old_rId)
            xml_slides.remove(old_el)
            new_el = xml_slides[-1]
            xml_slides.remove(new_el)
            xml_slides.insert(slide_index, new_el)
            affected = slide_index
        else:
            affected = len(prs.slides) - 1
        _OUT_DIR.mkdir(exist_ok=True)
        _save(prs, pptx_path)
        return _ok(slide_count=len(prs.slides), affected_index=affected, pptx_path=str(p))
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 7 — EXPORT
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def export_to_pdf(pptx_path: str, output_path: str = "") -> str:
    """Export a PPTX to PDF via LibreOffice (cross-platform) or comtypes (Windows).

    Args:
        pptx_path: Path to the source PPTX.
        output_path: Output PDF path. Defaults to same directory, .pdf extension.
    Returns: JSON: status ('ok' or 'unavailable'), pdf_path or message.
    """
    import subprocess, platform
    pptx_p = Path(pptx_path)
    if not pptx_p.exists(): return _err(f"File not found: {pptx_path}")
    out_p = Path(output_path) if output_path else pptx_p.with_suffix(".pdf")
    # Try LibreOffice
    for cmd in ["libreoffice", "soffice",
                "C:\\Program Files\\LibreOffice\\program\\soffice.exe"]:
        try:
            r = subprocess.run(
                [cmd, "--headless", "--convert-to", "pdf",
                 "--outdir", str(out_p.parent), str(pptx_p)],
                capture_output=True, timeout=90)
            if r.returncode == 0:
                auto = pptx_p.with_suffix(".pdf")
                if auto.exists() and auto != out_p: auto.rename(out_p)
                return _ok(pdf_path=str(out_p))
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    # Try comtypes on Windows
    if platform.system() == "Windows":
        try:
            import comtypes.client
            ppt = comtypes.client.CreateObject("Powerpoint.Application")
            ppt.Visible = 1
            deck = ppt.Presentations.Open(str(pptx_p.resolve()))
            deck.SaveAs(str(out_p.resolve()), 32)
            deck.Close(); ppt.Quit()
            return _ok(pdf_path=str(out_p))
        except Exception:
            pass
    return json.dumps({"status": "unavailable",
                       "message": "Install LibreOffice and ensure 'soffice' is on PATH.",
                       "pptx_path": str(pptx_p)}, ensure_ascii=False)


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 8 — ORIGINAL DESIGN RENDERING (unchanged API)
# ═══════════════════════════════════════════════════════════════════════════════

_JS_RENDERER = _RENDERERS / "pptxgenjs" / "index.js"


def _render_js(spec_path: Path, pptx_path: Path) -> str:
    """Run the PptxGenJS renderer via Node.js; returns stdout."""
    import subprocess, shutil
    node = shutil.which("node") or "node"
    result = subprocess.run(
        [node, str(_JS_RENDERER), str(spec_path), str(pptx_path)],
        capture_output=True, text=True, encoding="utf-8",
        cwd=str(_JS_RENDERER.parent),
    )
    if result.returncode != 0:
        raise RuntimeError(result.stderr or result.stdout or "node exited non-zero")
    return result.stdout


@mcp.tool()
def render_presentation(spec_json: str, filename: str = "", renderer: str = "js") -> str:
    """Render a complete JSON presentation spec to a PPTX file (PptxGenJS renderer by default).

    Call list_layouts() first to see the schema for each layout type.
    Call get_style_guide() for layout selection tips and color vocabulary.
    Call list_themes() to see available theme presets.

    Args:
        spec_json: Full presentation spec as a JSON string.
            Structure: {"meta": {"title": "...", "theme": "tech_dark", "colors": {}}, "slides": [...]}
        filename: Output filename without extension. Defaults to meta.title.
        renderer: "js" (default, PptxGenJS — better visuals) or "python" (legacy engine).
    Returns:
        JSON: status, pptx_path, slide_count, message.
    """
    try:
        spec = json.loads(spec_json)
    except json.JSONDecodeError as e:
        return json.dumps({"status": "error", "message": f"Invalid JSON: {e}"})
    _OUT_DIR.mkdir(exist_ok=True)
    if not filename:
        raw = spec.get("meta", {}).get("title", "presentation")
        filename = "".join(c if c.isalnum() or c in " -_" else "_" for c in raw).strip() or "presentation"
    spec_path = _OUT_DIR / f"{filename}.json"
    pptx_path = _OUT_DIR / f"{filename}.pptx"
    spec_path.write_text(json.dumps(spec, ensure_ascii=False, indent=2), encoding="utf-8")
    try:
        if renderer == "python":
            import engine
            engine.render(str(spec_path), str(pptx_path))
        else:
            _render_js(spec_path, pptx_path)
    except Exception as e:
        return json.dumps({"status": "error", "message": str(e)})
    return json.dumps({
        "status": "ok", "pptx_path": str(pptx_path),
        "slide_count": len(spec.get("slides", [])),
        "message": f"Generated {len(spec.get('slides',[]))} slides → {pptx_path.name}",
    }, ensure_ascii=False)


@mcp.tool()
def list_layouts() -> str:
    """Return the JSON schema for all 11 slide layouts.
    Always call this before generating a spec to know which fields each layout requires.
    """
    schema = {
        "cover":       {"desc": "封面 — 深色背景，大标题，装饰圆环", "required": ["title"],
                        "optional": {"tag": "顶部标签", "subtitle": "副标题", "subtitle2": "第二副标题",
                                     "badge": "底部徽章", "page": "页码 '01/08'"}},
        "hero_stat":   {"desc": "数据页 — 巨型数字 + 3列说明卡片",
                        "required": ["title", "hero_number", "hero_unit", "hero_label"],
                        "optional": {"columns": "最多3项，每项 {title, body}", "page": "页码"}},
        "solution":    {"desc": "解决方案 — 左深色/右白色分栏",
                        "required": ["hero", "bullets", "cards"],
                        "notes": "bullets最多3项{title,body}；cards最多3项{icon_char,title,body}",
                        "optional": {"tag": "顶部标签", "page": "页码"}},
        "pipeline":    {"desc": "流程 — 水平N阶段（最多5）带箭头", "required": ["title", "stages"],
                        "notes": "stages每项: {title, bullets:[最多4条], optional:bool}",
                        "optional": {"page": "页码"}},
        "hub_spoke":   {"desc": "中心辐射 — 中心概念 + 最多6个功能卡",
                        "required": ["title", "hub", "spokes"],
                        "notes": "hub:{title}；spokes每项:{title,body}，最多6项",
                        "optional": {"page": "页码"}},
        "kpi_grid":    {"desc": "KPI卡片 — 2×2网格", "required": ["title", "cards"],
                        "notes": "cards固定4项，每项:{value,unit,label,note?,color?,dark?}",
                        "optional": {"page": "页码"}},
        "radar":       {"desc": "雷达图 — 左侧图 + 右侧评分条",
                        "required": ["title", "dimensions"],
                        "notes": "dimensions每项:{name,us(0-10),avg(0-10)}，3-8项",
                        "optional": {"legend_us": "我方图例", "legend_avg": "对比图例", "page": "页码"}},
        "comparison":  {"desc": "竞品对比 — 功能矩阵表，我方高亮",
                        "required": ["title", "competitors", "features"],
                        "notes": "competitors最多4项{name,note}；features最多8项{name,values:[true/false/str]}",
                        "optional": {"page": "页码"}},
        "business":    {"desc": "商业模式 — 垂直时间轴 + 定价卡片",
                        "required": ["title", "items"],
                        "notes": "items最多3项:{title,body,price,price_label,color?}",
                        "optional": {"subtitle": "副标题", "page": "页码"}},
        "timeline":    {"desc": "时间轴 — 水平里程碑，上下交替卡片",
                        "required": ["title", "milestones"],
                        "notes": "milestones最多5项:{date,title,body,state:current|planned|target}",
                        "optional": {"footer": "底部注释", "page": "页码"}},
        "fundraising": {"desc": "融资页 — 深色背景，左侧金额，右侧资金用途",
                        "required": ["title", "amount", "breakdowns"],
                        "notes": "breakdowns最多3项:{pct,title,subtitle,body,color?}",
                        "optional": {"tag": "顶部标签", "amount_unit": "万", "round_label": "轮次",
                                     "target_time": "完成时间", "target_arr": "ARR目标",
                                     "contact": "联系方式", "page": "页码"}},
    }
    return json.dumps(schema, ensure_ascii=False, indent=2)


@mcp.tool()
def get_style_guide() -> str:
    """Return style vocabulary, layout selection guide, color palette, and spec-writing rules.
    Call when deciding which layout to use or how to structure a deck.
    """
    guide = {
        "layout_by_intent": {
            "封面/开场": "cover", "一个核心数字+说明": "hero_stat",
            "解决方案/产品": "solution", "流程/步骤": "pipeline",
            "中心概念+多功能": "hub_spoke", "4个KPI指标": "kpi_grid",
            "多维雷达对比": "radar", "竞品功能对比": "comparison",
            "商业模式/定价": "business", "路线图": "timeline", "融资/CTA": "fundraising",
        },
        "deck_templates": {
            "6页快速路演":  ["cover","hero_stat","solution","kpi_grid","business","fundraising"],
            "8页标准路演":  ["cover","hero_stat","solution","hub_spoke","kpi_grid","comparison","business","fundraising"],
            "11页完整路演": ["cover","hero_stat","solution","pipeline","hub_spoke","kpi_grid","radar","comparison","business","timeline","fundraising"],
        },
        "color_palette": {
            "accent": "#0EA5E9 (Sky Blue)", "green": "#22C55E (positive)",
            "indigo": "#6366F1 (premium)", "amber": "#F59E0B (highlight)",
            "dark_bg": "#1E293B (dark slides)",
        },
        "animation_types": ["appear","fade","fly","cut","zoom","wipe","split","blinds",
                            "checkerboard","dissolve","random_bars","peek","wheel","box",
                            "circle","diamond","plus","strips","wedge","stretch","expand",
                            "swivel","mixed","random"],
        "transition_types": ["fade","push","wipe","split","strips","cover","random"],
        "rules": [
            "page field format: '01/08' (denominator = total slide count)",
            "comparison.features.values order must match competitors order",
            "hero_stat.hero_number = digits only; unit goes in hero_unit",
            "timeline.milestones: only one state=current allowed",
            "kpi_grid: exactly 4 cards; card.dark=true uses dark background",
            "render_slide() can append/replace individual slides in existing decks",
            "add_shape_animation() shape_indices are zero-based shape list positions",
        ],
    }
    return json.dumps(guide, ensure_ascii=False, indent=2)



# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 9 — SHAPE POSITION / RESIZE
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def update_shape_position(pptx_path: str, slide_index: int, shape_index: int,
                          x: float = -1, y: float = -1,
                          w: float = -1, h: float = -1) -> str:
    """Move or resize an existing shape. Pass -1 to keep the current value.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        shape_index: Zero-based shape index on that slide.
        x, y: New top-left position in pixels (96 DPI). -1 = unchanged.
        w, h: New width / height in pixels. -1 = unchanged.
    Returns: JSON: status, x, y, w, h (final pixel values).
    """
    try:
        prs = _load(pptx_path)
        shape = prs.slides[slide_index].shapes[shape_index]
        if x >= 0: shape.left   = _emu(x)
        if y >= 0: shape.top    = _emu(y)
        if w >= 0: shape.width  = _emu(w)
        if h >= 0: shape.height = _emu(h)
        _save(prs, pptx_path)
        return _ok(x=_pxv(shape.left), y=_pxv(shape.top),
                   w=_pxv(shape.width), h=_pxv(shape.height))
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 10 — DATA-DRIVEN CHARTS
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def add_chart(pptx_path: str, slide_index: int, chart_type: str,
              categories: list, series: list,
              x: float, y: float, w: float, h: float,
              title: str = "", has_legend: bool = True) -> str:
    """Insert a data-driven, natively editable chart onto a slide.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        chart_type: 'column' | 'column_stacked' | 'bar' | 'bar_stacked' |
                    'line' | 'line_markers' | 'pie' | 'doughnut' | 'area'.
        categories: Category label list e.g. ['Q1', 'Q2', 'Q3', 'Q4'].
        series: List of series dicts: [{"name": "Revenue", "values": [4.3, 2.5, 3.5, 4.5]}].
        x, y, w, h: Position and size in pixels.
        title: Chart title text. '' = no title.
        has_legend: Show legend (default True).
    Returns: JSON: status, shape_index.
    """
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    _CT = {
        "column":         XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar":            XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_stacked":    XL_CHART_TYPE.BAR_STACKED,
        "line":           XL_CHART_TYPE.LINE,
        "line_markers":   XL_CHART_TYPE.LINE_MARKERS,
        "pie":            XL_CHART_TYPE.PIE,
        "doughnut":       XL_CHART_TYPE.DOUGHNUT,
        "area":           XL_CHART_TYPE.AREA,
    }
    try:
        prs = _load(pptx_path)
        slide = prs.slides[slide_index]
        cd = ChartData()
        cd.categories = categories
        for s in series:
            cd.add_series(s["name"], tuple(s["values"]))
        gf = slide.shapes.add_chart(
            _CT.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED),
            _emu(x), _emu(y), _emu(w), _emu(h), cd)
        chart = gf.chart
        chart.has_title = bool(title)
        if title: chart.chart_title.text_frame.text = title
        chart.has_legend = has_legend
        _save(prs, pptx_path)
        return _ok(shape_index=len(slide.shapes) - 1)
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 11 — THEME COLORS
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def set_theme_colors(pptx_path: str, accent1: str = "", accent2: str = "",
                     accent3: str = "", dark1: str = "", light1: str = "") -> str:
    """Modify the Office theme color scheme in the presentation.

    Changes affect built-in theme color slots visible in PowerPoint's color picker.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        accent1, accent2, accent3: Hex color for accent slots (primary brand colors).
        dark1: Hex for main dark/text color.
        light1: Hex for main light/background color.
    Returns: JSON: status, changed (list of slot names updated).
    """
    try:
        from lxml import etree
        THEME_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        prs = _load(pptx_path)
        theme_part = next(
            (rel.target_part for rel in prs.slide_masters[0].part.rels.values()
             if rel.reltype == THEME_RT), None)
        if theme_part is None:
            return _err("Theme part not found")
        theme_el = etree.fromstring(theme_part.blob)
        clr = theme_el.find(f".//{{{_A}}}clrScheme")
        if clr is None:
            return _err("Color scheme not found in theme")
        slots = {"accent1": accent1, "accent2": accent2, "accent3": accent3,
                 "dk1": dark1, "lt1": light1}
        changed = []
        for slot, val in slots.items():
            if not val: continue
            el = clr.find(f"{{{_A}}}{slot}")
            if el is None: continue
            for ch in list(el): el.remove(ch)
            etree.SubElement(el, f"{{{_A}}}srgbClr", val=val.lstrip("#"))
            changed.append(slot)
        theme_part.blob = etree.tostring(
            theme_el, xml_declaration=True, encoding="UTF-8", standalone=True)
        _save(prs, pptx_path)
        return _ok(changed=changed)
    except Exception as e: return _err(str(e))


# ═══════════════════════════════════════════════════════════════════════════════
# CATEGORY 12 — AUDIO / SOUND
# ═══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def set_slide_audio(pptx_path: str, slide_index: int, audio_path: str,
                    play_mode: str = "auto", loop: bool = False) -> str:
    """Embed an audio file on a slide.

    Supports MP3, WAV, M4A, OGG, AAC. The audio icon is placed at the
    bottom-right corner (invisible during presentation in auto mode).
    Note: if the slide already has shape animations, they will be replaced
    when play_mode='auto'. Use play_mode='click' to avoid this.

    Args:
        pptx_path: Path to PPTX. Modified in place.
        slide_index: Zero-based slide index.
        audio_path: Absolute path to the audio file.
        play_mode: 'auto' = play immediately on slide entry (default),
                   'click' = play only when clicked.
        loop: Loop the audio until the next slide (default False).
    Returns: JSON: status, shape_id, rId.
    """
    try:
        from lxml import etree
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI

        audio_p = Path(audio_path)
        if not audio_p.exists():
            return _err(f"Audio file not found: {audio_path}")
        audio_bytes = audio_p.read_bytes()
        suffix = audio_p.suffix.lower()
        _CT2 = {".mp3": "audio/mpeg", ".wav": "audio/wav",
                ".m4a": "audio/mp4",  ".ogg": "audio/ogg", ".aac": "audio/aac"}
        content_type = _CT2.get(suffix, "audio/mpeg")
        AUDIO_RT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/audio"

        prs = _load(pptx_path)
        if not (0 <= slide_index < len(prs.slides)):
            return _err(f"slide_index {slide_index} out of range")
        slide = prs.slides[slide_index]

        # Embed audio as a media part and add relationship
        media_n = len(slide.part.rels) + 1
        audio_uri = PackURI(f"/ppt/media/audio_{slide_index}_{media_n}{suffix}")
        audio_part = Part(audio_uri, content_type, audio_bytes)
        rId = slide.part.relate_to(audio_part, AUDIO_RT)

        # Unique shape ID
        shape_id = max((s.shape_id for s in slide.shapes), default=1) + 1

        # Audio shape (hidden icon at bottom-right)
        sp_xml = (
            f'<p:sp xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}">'
            f'<p:nvSpPr>'
            f'<p:cNvPr id="{shape_id}" name="Audio {shape_id}" descr="{audio_p.name}">'
            f'<a:hlinkClick r:id="{rId}" action="ppaction://media"/></p:cNvPr>'
            f'<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
            f'<p:nvPr/>'
            f'</p:nvSpPr>'
            f'<p:spPr>'
            f'<a:xfrm><a:off x="11176050" y="6337275"/>'
            f'<a:ext cx="457200" cy="457200"/></a:xfrm>'
            f'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom>'
            f'<a:ln><a:noFill/></a:ln>'
            f'</p:spPr>'
            f'<p:txBody><a:bodyPr/><a:lstStyle/>'
            f'<a:p><a:endParaRPr lang="en-US" dirty="0"/></a:p>'
            f'</p:txBody></p:sp>'
        )
        slide.shapes._spTree.append(etree.fromstring(sp_xml.encode("utf-8")))

        if play_mode == "auto":
            loop_attr = ' repeatDur="indefinite"' if loop else ""
            timing_xml = (
                f'<p:timing xmlns:p="{_P}">'
                f'<p:tnLst><p:par>'
                f'<p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">'
                f'<p:childTnLst>'
                f'<p:seq concurrent="1" nextAc="seek">'
                f'<p:cTn id="2" dur="indefinite" nodeType="mainSeq">'
                f'<p:childTnLst/></p:cTn>'
                f'<p:prevCondLst><p:cond evt="onPrev" delay="0">'
                f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>'
                f'<p:nextCondLst><p:cond evt="onNext" delay="0">'
                f'<p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>'
                f'</p:seq>'
                f'<p:audio isNarration="0">'
                f'<p:cMediaNode restart="never" fill="hold" syncBehavior="locked"{loop_attr}>'
                f'<p:cTn id="3" masterRel="sameClick" fill="hold">'
                f'<p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>'
                f'<p:tgtEl><p:spTgt spid="{shape_id}"/></p:tgtEl>'
                f'</p:cMediaNode></p:audio>'
                f'</p:childTnLst></p:cTn>'
                f'</p:par></p:tnLst>'
                f'</p:timing>'
            )
            sp_el = slide._element
            for old in sp_el.findall(f"{{{_P}}}timing"): sp_el.remove(old)
            sp_el.append(etree.fromstring(timing_xml.encode("utf-8")))

        _save(prs, pptx_path)
        return _ok(shape_id=shape_id, rId=rId)
    except Exception as e: return _err(str(e))


# ══════════════════════════════════════════════════════════════════════════════
# Category 13 — Theme / Template / Batch
# ══════════════════════════════════════════════════════════════════════════════

@mcp.tool()
def list_themes() -> dict:
    """
    List all available built-in theme presets.
    Returns names and brief descriptions so the AI can pick the right style.
    """
    try:
        import themes as _t
        info = {
            "tech_dark":  "深色科技风 — 深蓝底+蓝青色，适合AI/互联网/科技主题（默认）",
            "corporate":  "商务简洁风 — 白底+海军蓝，无光晕，适合金融/咨询/企业汇报",
            "warm":       "暖色活力风 — 奶油底+琥珀橙，适合消费品/零售/创业路演",
            "purple":     "紫霓科创风 — 深紫底+紫粉色，适合AI产品/创意/设计行业",
            "forest":     "自然生态风 — 绿色系，适合可持续/健康/农业/环保主题",
            "minimal":    "极简商务风 — 纯白+炭黑，无装饰，适合高端品牌/咨询",
            "sunset":     "落日创意风 — 暗暖底+橙红色，适合文创/广告/品牌策略",
            "ocean":      "深海数据风 — 深青底+蓝绿色，适合数据分析/BI/医疗健康",
        }
        return _ok(themes=list(_t.THEME_NAMES), descriptions=info)
    except Exception as e: return _err(str(e))


@mcp.tool()
def save_template(name: str, deck_spec: dict, description: str = "") -> dict:
    """
    Save a deck spec as a reusable template.
    Supports {{variable}} placeholders in any string value for batch rendering.

    Parameters
    ----------
    name        : template identifier (letters/digits/underscore/hyphen)
    deck_spec   : full deck JSON dict — same format as render_presentation
    description : short human-readable description

    Example deck_spec with placeholders:
        {
          "meta": {"theme": "corporate"},
          "slides": [
            {"layout": "cover", "title": "{{company}} 年度报告",
             "subtitle": "{{year}}年", "page": "01"},
            {"layout": "hero_stat", "hero_number": "{{revenue}}",
             "hero_unit": "{{unit}}", "title": "营收概览"}
          ]
        }
    """
    try:
        import re
        if not re.match(r'^[\w\-]+$', name):
            return _err("name must contain only letters, digits, _ or -")
        tmpl = {
            "name":        name,
            "description": description,
            "created":     __import__("datetime").date.today().isoformat(),
            "deck":        deck_spec,
        }
        path = _TMPL_DIR / f"{name}.json"
        path.write_text(json.dumps(tmpl, ensure_ascii=False, indent=2), encoding="utf-8")
        return _ok(saved=str(path), slide_count=len(deck_spec.get("slides", [])),
                   theme=deck_spec.get("meta", {}).get("theme", "tech_dark"))
    except Exception as e: return _err(str(e))


@mcp.tool()
def list_templates() -> dict:
    """List all saved templates with name, description, theme, slide count."""
    try:
        result = []
        for p in sorted(_TMPL_DIR.glob("*.json")):
            try:
                t = json.loads(p.read_text(encoding="utf-8"))
                result.append({
                    "name":        t.get("name", p.stem),
                    "description": t.get("description", ""),
                    "theme":       t.get("deck", {}).get("meta", {}).get("theme", "tech_dark"),
                    "slides":      len(t.get("deck", {}).get("slides", [])),
                    "created":     t.get("created", ""),
                })
            except Exception:
                pass
        return _ok(templates=result, count=len(result))
    except Exception as e: return _err(str(e))


@mcp.tool()
def render_from_template(template_name: str, variables: dict,
                          output_path: str = "") -> dict:
    """
    Render a saved template with variable substitution.

    Parameters
    ----------
    template_name : name used in save_template
    variables     : dict mapping placeholder names to replacement strings
                    e.g. {"company": "字节跳动", "year": "2026", "revenue": "3200"}
    output_path   : where to save the PPTX; omit to auto-name in projects/

    How substitution works
    ----------------------
    Every string value in the deck JSON is scanned for {{key}} patterns.
    Keys not found in variables are left as-is.
    """
    try:
        path = _TMPL_DIR / f"{template_name}.json"
        if not path.exists():
            return _err(f"Template '{template_name}' not found. Use list_templates() to see available templates.")
        tmpl    = json.loads(path.read_text(encoding="utf-8"))
        deck    = copy.deepcopy(tmpl["deck"])

        def _sub(obj):
            if isinstance(obj, str):
                for k, v in variables.items():
                    obj = obj.replace(f"{{{{{k}}}}}", str(v))
                return obj
            if isinstance(obj, dict):  return {k: _sub(v) for k, v in obj.items()}
            if isinstance(obj, list):  return [_sub(i) for i in obj]
            return obj

        deck = _sub(deck)

        if not output_path:
            safe = __import__("re").sub(r'[^\w\-]', '_',
                    variables.get("company", variables.get("name", template_name)))
            output_path = str(_OUT_DIR / f"{safe}.pptx")

        import tempfile
        tmp = Path(tempfile.mktemp(suffix=".json"))
        tmp.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
        try:
            _render_js(tmp, Path(output_path))
        except Exception:
            import engine as _engine
            _engine.render(str(tmp), output_path)
        finally:
            tmp.unlink(missing_ok=True)
        return _ok(output=output_path, template=template_name, variables=variables)
    except Exception as e: return _err(str(e))


@mcp.tool()
def batch_render(template_name: str, items: list,
                  output_dir: str = "", filename_key: str = "name") -> dict:
    """
    Render one PPTX per item in items[] using a saved template.
    Each item is a variables dict passed to render_from_template.

    Parameters
    ----------
    template_name : name used in save_template
    items         : list of variable dicts, one per output file
                    e.g. [{"company":"A公司","revenue":"200万"},
                          {"company":"B公司","revenue":"350万"}]
    output_dir    : directory for output files (default: projects/)
    filename_key  : which variable to use as filename (default: "name")

    Returns
    -------
    List of generated file paths and any per-item errors.
    """
    try:
        out_dir = Path(output_dir) if output_dir else _OUT_DIR
        out_dir.mkdir(parents=True, exist_ok=True)

        path = _TMPL_DIR / f"{template_name}.json"
        if not path.exists():
            return _err(f"Template '{template_name}' not found.")
        tmpl = json.loads(path.read_text(encoding="utf-8"))

        import tempfile, re

        def _sub(obj, variables):
            if isinstance(obj, str):
                for k, v in variables.items():
                    obj = obj.replace(f"{{{{{k}}}}}", str(v))
                return obj
            if isinstance(obj, dict):  return {k: _sub(v, variables) for k, v in obj.items()}
            if isinstance(obj, list):  return [_sub(i, variables) for i in obj]
            return obj

        results, errors = [], []
        for i, item in enumerate(items):
            try:
                deck     = _sub(copy.deepcopy(tmpl["deck"]), item)
                raw_name = str(item.get(filename_key, f"item_{i+1:03d}"))
                safe     = re.sub(r'[^\w\-]', '_', raw_name)
                out_path = str(out_dir / f"{safe}.pptx")
                tmp      = Path(tempfile.mktemp(suffix=".json"))
                tmp.write_text(json.dumps(deck, ensure_ascii=False), encoding="utf-8")
                try:
                    _render_js(tmp, Path(out_path))
                except Exception:
                    import engine as _engine
                    _engine.render(str(tmp), out_path)
                finally:
                    tmp.unlink(missing_ok=True)
                results.append(out_path)
            except Exception as ex:
                errors.append({"index": i, "item": item, "error": str(ex)})

        return _ok(generated=len(results), files=results,
                   errors=errors, output_dir=str(out_dir))
    except Exception as e: return _err(str(e))


# ── Entry point ─────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    mcp.run()


