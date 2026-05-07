#!/usr/bin/env python3
"""
JSON → PPTX renderer.
Usage:  python engine.py slides.json output.pptx
"""
import sys
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from lxml import etree

sys.path.insert(0, str(Path(__file__).parent))
sys.path.insert(0, str(Path(__file__).parent.parent))

_P   = "http://schemas.openxmlformats.org/presentationml/2006/main"
_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_A   = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"


def _pptx_el(xml_str: str):
    """Parse a PPTX XML fragment that uses p:/a:/r: prefixes."""
    wrapped = (f'<root xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}"'
               f' xmlns:p14="{_P14}">{xml_str.strip()}</root>')
    return etree.fromstring(wrapped.encode("utf-8"))[0]


def _apply_slide_effects(slide, spec: dict):
    """Inject transition and/or animation declared inline in the slide spec."""
    # ── Transition ────────────────────────────────────────────────────────────
    trans = spec.get("transition")
    if trans:
        try:
            from pptx_animations import create_transition_xml
            xml = create_transition_xml(
                trans.get("type", "fade"),
                duration=trans.get("duration_ms", 500) / 1000,
                advance_after=(trans.get("advance_after_ms", 0) / 1000) or None,
            )
            sp_el = slide._element
            for old in sp_el.findall(f"{{{_P}}}transition"): sp_el.remove(old)
            sp_el.append(_pptx_el(xml))
        except Exception as exc:
            print(f"[WARN] transition failed: {exc}")

    # ── Entrance animation ────────────────────────────────────────────────────
    anim = spec.get("animation")
    if anim:
        try:
            from pptx_animations import create_sequence_timing_xml, pick_animation_effect
            shape_indices = anim.get("shapes", [])
            anim_type  = anim.get("type", "fade")
            trigger    = anim.get("trigger", "after-previous")
            dur_ms     = anim.get("duration_ms", 500)
            delay_ms   = anim.get("delay_ms", 0)
            targets = [
                (slide.shapes[idx].shape_id, delay_ms,
                 pick_animation_effect(anim_type, i))
                for i, idx in enumerate(shape_indices)
                if 0 <= idx < len(slide.shapes)
            ]
            if targets:
                xml = create_sequence_timing_xml(
                    targets, duration=dur_ms / 1000, trigger=trigger)
                sp_el = slide._element
                for old in sp_el.findall(f"{{{_P}}}timing"): sp_el.remove(old)
                sp_el.append(_pptx_el(xml))
        except Exception as exc:
            print(f"[WARN] animation failed: {exc}")

import theme
import layouts
import themes as _themes_mod

DISPATCH = {
    "cover":       layouts.render_cover,
    "hero_stat":   layouts.render_hero_stat,
    "solution":    layouts.render_solution,
    "pipeline":    layouts.render_pipeline,
    "hub_spoke":   layouts.render_hub_spoke,
    "kpi_grid":    layouts.render_kpi_grid,
    "radar":       layouts.render_radar,
    "comparison":  layouts.render_comparison,
    "business":    layouts.render_business,
    "timeline":    layouts.render_timeline,
    "fundraising": layouts.render_fundraising,
}


def _apply_colors(colors: dict) -> dict:
    """Override theme constants; return originals for later restore."""
    originals = {}
    for k, v in colors.items():
        if hasattr(theme, k):
            originals[k] = getattr(theme, k)
            setattr(theme, k, tuple(v) if isinstance(v, list) else v)
    return originals


def _restore_colors(originals: dict):
    for k, v in originals.items():
        setattr(theme, k, v)


def _apply_theme(theme_name: str, extra_colors: dict) -> dict:
    """
    Apply a named preset then individual overrides in one pass.
    Returns the combined originals dict for later restore.
    """
    preset = _themes_mod.get_theme(theme_name)
    combined = {**preset, **extra_colors}   # extra_colors win over preset
    return _apply_colors(combined)


def render(spec_path: str, out_path: str):
    with open(spec_path, encoding="utf-8") as f:
        deck = json.load(f)

    meta      = deck.get("meta", {})
    originals = _apply_theme(meta.get("theme", "tech_dark"),
                             meta.get("colors", {}))

    prs = Presentation()
    prs.slide_width  = Emu(theme.SLIDE_W_PX * theme.PX_TO_EMU)
    prs.slide_height = Emu(theme.SLIDE_H_PX * theme.PX_TO_EMU)
    blank = prs.slide_layouts[6]

    for i, spec in enumerate(deck["slides"], 1):
        fn = DISPATCH.get(spec["layout"])
        if fn is None:
            print(f"[SKIP] slide {i}: unknown layout '{spec['layout']}'")
            continue
        slide = prs.slides.add_slide(blank)
        fn(slide, spec)
        _apply_slide_effects(slide, spec)
        print(f"[OK]   {i:02d}  {spec.get('id', spec['layout'])}")

    prs.save(out_path)
    _restore_colors(originals)
    print(f"\nSaved → {out_path}")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python engine.py <slides.json> <output.pptx>")
        sys.exit(1)
    render(sys.argv[1], sys.argv[2])
