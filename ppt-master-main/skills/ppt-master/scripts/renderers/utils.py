from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import theme

# ── Coordinate ────────────────────────────────────────────────────────────────

def px(n: float) -> Emu:
    return Emu(int(n * theme.PX_TO_EMU))

# ── Background ────────────────────────────────────────────────────────────────

def set_bg(slide, hex_color: str):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = theme.rgb(hex_color)

# ── Shapes ────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=1.0, radius=0):
    shape = slide.shapes.add_shape(1, px(x), px(y), px(w), px(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = theme.rgb(line)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    if radius > 0:
        _apply_radius(shape, h, radius)
    return shape


def add_oval(slide, cx, cy, r, *, fill=None, line=None, line_w=1.5):
    shape = slide.shapes.add_shape(9, px(cx - r), px(cy - r), px(r * 2), px(r * 2))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = theme.rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = theme.rgb(line)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_line(slide, x1, y1, x2, y2, *, color=None, width=1.0):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    conn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                                      px(x1), px(y1), px(x2), px(y2))
    conn.line.color.rgb = theme.rgb(color or theme.BORDER)
    conn.line.width = Pt(width)
    return conn


def _apply_radius(shape, h_px: float, radius_px: float):
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    pg = spPr.find(qn("a:prstGeom"))
    if pg is None:
        return
    pg.set("prst", "roundRect")
    avLst = pg.find(qn("a:avLst"))
    if avLst is None:
        avLst = etree.SubElement(pg, qn("a:avLst"))
    for gd in avLst.findall(qn("a:gd")):
        avLst.remove(gd)
    gd = etree.SubElement(avLst, qn("a:gd"))
    gd.set("name", "adj")
    gd.set("fmla", f"val {min(int(radius_px / h_px * 100000), 50000)}")

# ── Text ──────────────────────────────────────────────────────────────────────

def add_text(slide, text: str, x, y, w, h, *,
             font=None, size=16, bold=False,
             color=theme.T_PRIMARY, align="left", wrap=True,
             line_spacing=None):
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT,
                   "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = theme.rgb(color)
    run.font.name = font or theme.FONT_CJK
    _set_ea_font(run, font or theme.FONT_CJK)
    return tb

# ── Common composites ─────────────────────────────────────────────────────────

def title_bar(slide, x, y, title, bar_w=80, size=32, color=theme.ACCENT):
    tc = "F1F5F9" if theme.is_dark_bg() else theme.T_PRIMARY
    add_text(slide, title, x, y, 1100, 52, size=size, bold=True, color=tc)
    add_rect(slide, x, y + int(size * 1.05), bar_w, 3, fill=color)


def page_num(slide, page: str):
    if page:
        tc = "475569" if theme.is_dark_bg() else theme.T_MUTED
        add_text(slide, page, 1180, 694, 80, 22, size=13, color=tc, align="right")


# ── Gradient helpers ──────────────────────────────────────────────────────────

def _grad_fill(c1: str, c2: str, angle_deg: int = 90):
    gf = etree.Element(qn("a:gradFill"))
    gsLst = etree.SubElement(gf, qn("a:gsLst"))
    for pos, c in (("0", c1.lstrip("#").upper()), ("100000", c2.lstrip("#").upper())):
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", pos)
        etree.SubElement(gs, qn("a:srgbClr")).set("val", c)
    lin = etree.SubElement(gf, qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000) % 21600000))
    lin.set("scaled", "0")
    return gf


def set_grad_bg(slide, c1: str, c2: str, angle: int = 90):
    """Apply linear gradient to slide background (replaces solid fill)."""
    sp_tree = slide._element
    bg_el = sp_tree.find(qn("p:bg"))
    if bg_el is None:
        bg_el = etree.Element(qn("p:bg"))
        sp_tree.insert(0, bg_el)
    bgPr = bg_el.find(qn("p:bgPr"))
    if bgPr is None:
        bgPr = etree.SubElement(bg_el, qn("p:bgPr"))
    for child in list(bgPr):
        bgPr.remove(child)
    bgPr.insert(0, _grad_fill(c1, c2, angle))


def apply_grad(shape, c1: str, c2: str, angle: int = 90):
    """Replace shape solid fill with linear gradient."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    for child in list(spPr):
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")):
            spPr.remove(child)
    spPr.insert(0, _grad_fill(c1, c2, angle))


def apply_shadow(shape, blur: int = 50800, dist: int = 25400, alpha: int = 20):
    """Add subtle downward outer shadow to a shape."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    eff = spPr.find(qn("a:effectLst"))
    if eff is None:
        eff = etree.SubElement(spPr, qn("a:effectLst"))
    shdw = etree.SubElement(eff, qn("a:outerShdw"))
    shdw.set("blurRad", str(blur))
    shdw.set("dist",    str(dist))
    shdw.set("dir",     "5400000")   # 90° = downward
    shdw.set("algn",    "ctr")
    shdw.set("rotWithShape", "0")
    srgb = etree.SubElement(shdw, qn("a:srgbClr"))
    srgb.set("val", "000000")
    etree.SubElement(srgb, qn("a:alpha")).set("val", str(alpha * 1000))


def add_rect_alpha(slide, x, y, w, h, fill: str, alpha: int):
    """Rectangle with semi-transparent fill (alpha 0=invisible … 100=opaque)."""
    shape = add_rect(slide, x, y, w, h, fill=fill)
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    sf = spPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            etree.SubElement(srgb, qn("a:alpha")).set("val", str(int(alpha * 1000)))
    return shape


def set_bg_image(slide, image_path: str):
    """Embed image as bottom-most background (behind all shapes)."""
    from pptx.util import Emu as _Emu
    pic = slide.shapes.add_picture(
        image_path,
        _Emu(0), _Emu(0),
        _Emu(theme.SLIDE_W_PX * theme.PX_TO_EMU),
        _Emu(theme.SLIDE_H_PX * theme.PX_TO_EMU),
    )
    spTree = slide.shapes._spTree
    pic_el = pic._element
    spTree.remove(pic_el)
    spTree.insert(2, pic_el)   # index 2 = after nvGrpSpPr + grpSpPr
    return pic


# ── East-Asian font fix ───────────────────────────────────────────────────────

def _set_ea_font(run, typeface: str):
    """Inject <a:ea> after <a:latin> so CJK chars use the specified typeface."""
    r_el = run._r
    rPr = r_el.find(qn("a:rPr"))
    if rPr is None:
        return
    for ea in rPr.findall(qn("a:ea")):
        rPr.remove(ea)
    ea_el = etree.Element(qn("a:ea"))
    ea_el.set("typeface", typeface)
    lat = rPr.find(qn("a:latin"))
    if lat is not None:
        lat.addnext(ea_el)
    else:
        rPr.append(ea_el)


# ── Radial gradient / ambient glow ───────────────────────────────────────────

def _radial_grad_fill(color: str, alpha_center: int = 30, alpha_edge: int = 0):
    """Radial gradient (circle) from colored centre to transparent edge."""
    gf = etree.Element(qn("a:gradFill"))
    gsLst = etree.SubElement(gf, qn("a:gsLst"))
    c = color.lstrip("#").upper()
    for pos, al in (("0", alpha_center), ("100000", alpha_edge)):
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", pos)
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", c)
        etree.SubElement(srgb, qn("a:alpha")).set("val", str(al * 1000))
    path_el = etree.SubElement(gf, qn("a:path"))
    path_el.set("path", "circle")
    ftr = etree.SubElement(path_el, qn("a:fillToRect"))
    ftr.set("l", "50000"); ftr.set("t", "50000")
    ftr.set("r", "50000"); ftr.set("b", "50000")
    return gf


def add_glow_orb(slide, cx: int, cy: int, r: int, color: str, alpha: int = 25):
    """Large ambient glow circle: radial gradient from colour to transparent."""
    if not theme.DECO_GLOW:
        return None
    shape = slide.shapes.add_shape(9, px(cx - r), px(cy - r), px(r * 2), px(r * 2))
    shape.fill.background()
    shape.line.fill.background()
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    for child in list(spPr):
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill"), qn("a:blipFill")):
            spPr.remove(child)
    spPr.insert(0, _radial_grad_fill(color, alpha, 0))
    return shape


# ── Multi-stop gradient ───────────────────────────────────────────────────────

def _multi_grad_fill(stops: list, angle_deg: int = 90):
    """Multi-stop linear gradient. stops = [(pos_pct_0-100, hex), ...]"""
    gf = etree.Element(qn("a:gradFill"))
    gsLst = etree.SubElement(gf, qn("a:gsLst"))
    for pos_pct, c in stops:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(int(pos_pct * 1000)))
        etree.SubElement(gs, qn("a:srgbClr")).set("val", c.lstrip("#").upper())
    lin = etree.SubElement(gf, qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000) % 21600000))
    lin.set("scaled", "0")
    return gf


def set_multi_grad_bg(slide, stops: list, angle: int = 90):
    """Apply multi-stop gradient to slide background."""
    sp_tree = slide._element
    bg_el = sp_tree.find(qn("p:bg"))
    if bg_el is None:
        bg_el = etree.Element(qn("p:bg"))
        sp_tree.insert(0, bg_el)
    bgPr = bg_el.find(qn("p:bgPr"))
    if bgPr is None:
        bgPr = etree.SubElement(bg_el, qn("p:bgPr"))
    for child in list(bgPr):
        bgPr.remove(child)
    bgPr.insert(0, _multi_grad_fill(stops, angle))


# ── Glassmorphism ─────────────────────────────────────────────────────────────

def apply_glass(shape, fill_alpha: int = 10):
    """Make shape fill semi-transparent for glassmorphism. fill_alpha: 0-100."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    sf = spPr.find(qn("a:solidFill"))
    if sf is not None:
        srgb = sf.find(qn("a:srgbClr"))
        if srgb is not None:
            for a in srgb.findall(qn("a:alpha")):
                srgb.remove(a)
            etree.SubElement(srgb, qn("a:alpha")).set("val", str(fill_alpha * 1000))


def apply_glass_border(shape, border_alpha: int = 20, border_color: str = "FFFFFF"):
    """Set shape border to semi-transparent (glass edge shimmer)."""
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    if spPr is None:
        return
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return
    for child in list(ln):
        if child.tag in (qn("a:solidFill"), qn("a:gradFill"), qn("a:noFill")):
            ln.remove(child)
    new_sf = etree.SubElement(ln, qn("a:solidFill"))
    srgb = etree.SubElement(new_sf, qn("a:srgbClr"))
    srgb.set("val", border_color.lstrip("#").upper())
    etree.SubElement(srgb, qn("a:alpha")).set("val", str(border_alpha * 1000))


# ── Noise / grain texture ─────────────────────────────────────────────────────

def _noise_png(alpha: int = 8) -> bytes:
    """Generate a grayscale+alpha noise PNG (pure stdlib, deterministic)."""
    import struct, zlib, random

    try:
        import numpy as np
        w, h = 320, 180
        rng = np.random.default_rng(42)
        gray = rng.integers(0, 256, (h, w), dtype=np.uint8)
        alpha_arr = np.full((h, w), alpha, dtype=np.uint8)
        pixels = np.stack([gray, alpha_arr], axis=-1)           # (h, w, 2)
        filter_col = np.zeros((h, 1), dtype=np.uint8)
        rows_data = np.concatenate([filter_col, pixels.reshape(h, w * 2)], axis=1)
        raw = rows_data.tobytes()
    except ImportError:
        w, h = 128, 72
        rng = random.Random(42)
        rows = []
        for _ in range(h):
            row = bytearray([0])
            for _ in range(w):
                row.extend([rng.randint(0, 255), alpha])
            rows.append(bytes(row))
        raw = b''.join(rows)

    compressed = zlib.compress(raw, 6)

    def mk_chunk(t, d):
        return struct.pack('>I', len(d)) + t + d + struct.pack('>I', zlib.crc32(t + d) & 0xFFFFFFFF)

    ihdr = struct.pack('>IIBBBBB', w, h, 8, 4, 0, 0, 0)   # colour-type 4 = grayscale+alpha
    return (b'\x89PNG\r\n\x1a\n'
            + mk_chunk(b'IHDR', ihdr)
            + mk_chunk(b'IDAT', compressed)
            + mk_chunk(b'IEND', b''))


def add_noise_texture(slide, alpha: int = 5):
    """Overlay a full-slide grain texture above background, below content."""
    if not theme.DECO_NOISE:
        return None
    import io
    from pptx.util import Emu as _Emu
    png = _noise_png(alpha)
    pic = slide.shapes.add_picture(
        io.BytesIO(png),
        _Emu(0), _Emu(0),
        _Emu(theme.SLIDE_W_PX * theme.PX_TO_EMU),
        _Emu(theme.SLIDE_H_PX * theme.PX_TO_EMU),
    )
    return pic
